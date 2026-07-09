import os
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from io import BytesIO
import pymupdf
import fitz
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from src.pdfsummarizer.logger import logging
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Extracting text from PDF/PPTX using OCR (PyMuPDF + pytesseract, no poppler dependency)
# and as well as making chunks of it.


def _iter_all_shapes(shapes):
    """Yield every shape on a slide, descending into group shapes.

    `slide.shapes` only iterates top-level shapes - anything a user has
    grouped (pictures, textboxes, tables) is otherwise skipped entirely.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)


def _extract_image_blob(shape):
    """Return raw image bytes if this shape is carrying a picture.

    Handles both plain Picture shapes (shape_type == PICTURE) AND picture
    placeholders (shape_type == PLACEHOLDER) - a placeholder with an image
    inserted into it still reports as PLACEHOLDER, not PICTURE, so checking
    only for PICTURE silently misses it.
    """
    try:
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER):
            return shape.image.blob
    except AttributeError:
        # Placeholder exists but has no picture inserted into it (e.g. an
        # empty title/body placeholder) - nothing to OCR.
        return None
    except NotImplementedError:
        # A handful of custom/freeform autoshapes don't implement
        # shape_type at all - treat as "not a picture" rather than crash.
        return None
    return None


def _extract_text_from_single_file(file):
    """Extract text from a single file (PDF, PPTX, TXT, MD, or CSV)."""
    file_extension = Path(file.name).suffix.lower()
    text = ""

    # PDF Files (OCR-based extraction via PyMuPDF, no poppler needed)
    if file_extension == ".pdf":
        file_bytes = file.read()
        pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")

        for page_number, page in enumerate(pdf_doc, start=1):
            pix = page.get_pixmap(dpi=300)
            image = Image.open(BytesIO(pix.tobytes("png")))
            page_text = pytesseract.image_to_string(image)
            text += (page_text or "") + "\n"
            logging.info(f"OCR completed for PDF page {page_number} from {file.name}.")

        pdf_doc.close()

    # PowerPoint Files (text extraction + OCR on embedded images)
    elif file_extension in [".ppt", ".pptx"]:
        if file_extension == ".ppt":
            # python-pptx only reads the modern OOXML .pptx format.
            # A real legacy binary .ppt fails with a cryptic
            # PackageNotFoundError, so fail loudly and clearly instead.
            raise ValueError(
                "Legacy .ppt files are not supported. python-pptx can "
                "only read the modern .pptx format - please save/export "
                "the file as .pptx and re-upload."
            )

        file_bytes = file.read()
        presentation = Presentation(BytesIO(file_bytes))

        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in _iter_all_shapes(slide.shapes):
                try:
                    # Native text (text boxes, titles, body placeholders, etc.)
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs)
                            if line:
                                text += line + "\n"

                    # Tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            text += " ".join(cell.text for cell in row.cells) + "\n"

                    # Pictures - both plain Picture shapes and picture
                    # placeholders that have an image inserted into them.
                    image_blob = _extract_image_blob(shape)
                    if image_blob:
                        try:
                            image = Image.open(BytesIO(image_blob))
                            ocr_text = pytesseract.image_to_string(image)
                            text += (ocr_text or "") + "\n"
                        except Exception as img_err:
                            logging.warning(
                                f"Could not OCR image on slide {slide_number} from {file.name}: {img_err}"
                            )

                except NotImplementedError:
                    # Some custom/freeform autoshapes don't implement
                    # shape_type / related properties - skip rather than
                    # crash the whole file over one odd shape.
                    logging.warning(
                        f"Skipped an unrecognized shape on slide {slide_number} from {file.name}."
                    )
                    continue

            logging.info(f"Processed slide {slide_number} from {file.name}.")

    # Plain Text Files
    elif file_extension in [".txt", ".md", ".csv"]:
        text = file.read().decode("utf-8")

    else:
        raise ValueError(
            f"Unsupported file type: {file_extension}. "
            "Supported formats are PDF, PPTX, TXT, MD, and CSV."
        )

    if not text.strip():
        raise ValueError(f"The file {file.name} contains no readable text.")

    return text


def extract_text_from_file(file):
    """Extract and chunk text from a single file."""
    try:
        text = _extract_text_from_single_file(file)
        
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=700,
            chunk_overlap=100
        )

        chunks = text_splitter.split_text(text)
        logging.info(f"Extracted {len(chunks)} chunks successfully from {file.name}.")
        return chunks

    except Exception as e:
        logging.exception(e)
        raise


def _extract_text_from_file_with_reset(file):
    """Extract text from a file after resetting its pointer to the start."""
    try:
        file.seek(0)
    except Exception:
        pass
    return _extract_text_from_single_file(file)


def extract_text_from_files(files):
    """
    Extract and chunk text from multiple files in parallel.

    Args:
        files: A list of file objects or a single file object

    Returns:
        A list of text chunks from all processed files
    """
    try:
        # Handle single file for backward compatibility
        if not isinstance(files, list):
            files = [files]

        if not files:
            raise ValueError("No files provided.")

        if len(files) == 1:
            file = files[0]
            logging.info(f"Processing file: {file.name}")
            text = _extract_text_from_file_with_reset(file)
            logging.info(f"Successfully extracted text from {file.name}.")
            all_text = text + "\n"
        else:
            max_workers = min(len(files), max(2, os.cpu_count() or 4))
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                future_map = {
                    executor.submit(_extract_text_from_file_with_reset, file): file
                    for file in files
                }

                extracted_texts = []
                for future in future_map:
                    file = future_map[future]
                    try:
                        logging.info(f"Processing file: {file.name}")
                        text = future.result()
                        extracted_texts.append(text)
                        logging.info(f"Successfully extracted text from {file.name}.")
                    except Exception as e:
                        logging.warning(f"Error processing {file.name}: {e}")

            all_text = "\n".join(text for text in extracted_texts if text)

        if not all_text.strip():
            raise ValueError("No readable text found in any of the provided files.")

        # Split combined text into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=700,
            chunk_overlap=100
        )

        chunks = text_splitter.split_text(all_text)
        logging.info(f"Extracted {len(chunks)} chunks successfully from {len(files)} file(s).")
        return chunks

    except Exception as e:
        logging.exception(e)
        raise